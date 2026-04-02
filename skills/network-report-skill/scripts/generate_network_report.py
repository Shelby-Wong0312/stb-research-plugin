#!/usr/bin/env python3
"""
人脈網絡報告 PPT Generator
完全複刻 葉寅夫_人脈網絡報告.pptx 的 18 頁格式

pip install python-pptx --break-system-packages
"""

import json
import argparse
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ==============================================================================
# 色彩系統（完全匹配原稿）
# ==============================================================================
C_DEEP_BLUE = RGBColor(0x1A, 0x36, 0x5D)    # 標題欄、主標題
C_MID_BLUE = RGBColor(0x2C, 0x52, 0x82)      # 章節小標
C_GOLD = RGBColor(0xC6, 0xA0, 0x52)          # 卡片標題 accent
C_GREEN = RGBColor(0x38, 0xA1, 0x69)         # 正向
C_RED = RGBColor(0xC5, 0x30, 0x30)           # 負向
C_DARK_TEXT = RGBColor(0x2D, 0x37, 0x48)     # 正文
C_GRAY = RGBColor(0x71, 0x80, 0x96)          # 來源、頁碼
C_LIGHT_GRAY = RGBColor(0xA0, 0xAE, 0xC0)   # 副標題
C_BROWN = RGBColor(0x78, 0x35, 0x0F)         # 語錄
C_DARK_BROWN = RGBColor(0x92, 0x40, 0x0E)    # 哲學標題
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_BLACK = RGBColor(0x00, 0x00, 0x00)
C_HEADER_BG = RGBColor(0xE2, 0xE8, 0xF0)    # 表頭淺灰
C_CARD_BG = RGBColor(0x2D, 0x37, 0x48)       # 卡片背景深灰
C_QUOTE_BG = RGBColor(0xFE, 0xF3, 0xC7)      # 語錄背景淺黃

# 投影片尺寸
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# 字體大小
F_COVER_NAME = Pt(60)
F_COVER_TITLE = Pt(22)
F_COVER_TAG = Pt(16)
F_PAGE_TITLE = Pt(32)
F_SECTION = Pt(15)
F_BODY = Pt(12)
F_TABLE = Pt(12)
F_TABLE_HEADER = Pt(12)
F_SMALL = Pt(11)
F_SOURCE = Pt(11)
F_PAGE_NUM = Pt(12)
F_CARD_TITLE = Pt(15)
F_CARD_SUB = Pt(11)
F_CARD_PEOPLE = Pt(10)
F_MILESTONE = Pt(12)
F_QUOTE = Pt(12)
F_TOC = Pt(22)

# Layout constants
MARGIN_L = Inches(0.4)
CONTENT_W = Inches(12.5)
HEADER_H = Inches(0.9)


# ==============================================================================
# Helper functions
# ==============================================================================

def _add_header_bar(slide):
    """添加頁面頂部的淺灰色 header 背景條"""
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        SLIDE_W, HEADER_H
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFC)
    bar.line.fill.background()
    # 頂部深藍細線
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        SLIDE_W, Inches(0.008)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = C_DEEP_BLUE
    line.line.fill.background()


def _add_page_title(slide, text):
    """添加頁面標題（header bar 內）"""
    _add_header_bar(slide)
    tb = slide.shapes.add_textbox(MARGIN_L, Inches(0.13), Inches(11), Inches(0.6))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = F_PAGE_TITLE
    p.font.bold = True
    p.font.color.rgb = C_DEEP_BLUE


def _add_source_line(slide, source):
    """
    添加頁面底部來源說明。支援兩種格式：
    1. 純文字：source = "來源：經濟日報、天下雜誌"
    2. 結構化：source = [{"title":"文章標題","url":"https://..."},...]
       會生成帶超連結的來源列表
    """
    if not source:
        return
    if isinstance(source, str):
        tb = slide.shapes.add_textbox(MARGIN_L, Inches(7.0), Inches(11.5), Inches(0.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = source
        p.font.size = F_SOURCE
        p.font.color.rgb = C_GRAY
    elif isinstance(source, list):
        tb = slide.shapes.add_textbox(MARGIN_L, Inches(6.75), Inches(11.5), Inches(0.65))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run0 = p.add_run()
        run0.text = "資料來源："
        run0.font.size = F_SOURCE
        run0.font.color.rgb = C_GRAY
        run0.font.bold = True
        for idx, s in enumerate(source[:4]):  # 最多顯示4筆
            title = s.get("title", "")
            url = s.get("url", "")
            if idx > 0:
                p_new = tf.add_paragraph()
            else:
                p_new = p
            run = p_new.add_run()
            display = f"[{idx+1}] {title}"
            if url and not title:
                display = f"[{idx+1}] {url[:60]}..."
            run.text = ("  " if idx == 0 else "") + display
            run.font.size = Pt(9)
            run.font.color.rgb = C_GRAY
            if url:
                run.hyperlink.address = url


def _add_page_num(slide, num):
    """添加頁碼"""
    tb = slide.shapes.add_textbox(Inches(12.4), Inches(7.1), Inches(0.5), Inches(0.3))
    p = tb.text_frame.paragraphs[0]
    p.text = str(num)
    p.font.size = F_PAGE_NUM
    p.font.color.rgb = C_GRAY
    p.alignment = PP_ALIGN.RIGHT


def _add_section_title(slide, text, left, top, width=Inches(5.5)):
    """添加章節小標題"""
    tb = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    p = tb.text_frame.paragraphs[0]
    p.text = text
    p.font.size = F_SECTION
    p.font.bold = True
    p.font.color.rgb = C_MID_BLUE


def _make_table(slide, data, left, top, width, col_widths, row_h=Inches(0.4)):
    """
    通用表格：第一行為 header（淺灰底、置中）
    data = [['col1','col2',...], ['val1','val2',...], ...]
    col_widths = relative widths list
    """
    if not data:
        return
    rows, cols = len(data), len(data[0])
    shape = slide.shapes.add_table(rows, cols, left, top, width, row_h * rows)
    tbl = shape.table
    total_w = sum(col_widths)
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = int(width * w / total_w)
    for r in range(rows):
        tbl.rows[r].height = row_h
    for r, row_data in enumerate(data):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r, c)
            cell.text = str(val) if val else ""
            tf = cell.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.05)
            tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.02)
            tf.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in tf.paragraphs:
                p.font.size = F_TABLE if r > 0 else F_TABLE_HEADER
                p.font.bold = (r == 0)
                p.font.color.rgb = C_DARK_TEXT
                p.alignment = PP_ALIGN.CENTER if r == 0 else PP_ALIGN.LEFT
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_HEADER_BG


def _add_rounded_rect(slide, left, top, width, height, fill_color):
    """添加圓角矩形"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text,
                 font_size=F_BODY, bold=False, color=C_DARK_TEXT,
                 alignment=PP_ALIGN.LEFT, word_wrap=True):
    """通用 textbox"""
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = word_wrap
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = font_size
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tb


# ==============================================================================
# Slide generators
# ==============================================================================

def slide_cover(prs, d):
    """Slide 1: 封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # 深藍全屏背景
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = C_DEEP_BLUE
    bg.line.fill.background()
    # 頂部細裝飾線
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.008))
    line.fill.solid()
    line.fill.fore_color.rgb = C_GOLD
    line.line.fill.background()
    # 人名
    _add_textbox(slide, Inches(0.5), Inches(2.5), Inches(12.3), Inches(1.2),
                 d.get("target_name", ""), F_COVER_NAME, True, C_WHITE)
    # 職稱
    _add_textbox(slide, Inches(4.3), Inches(3.7), Inches(4.7), Inches(0.5),
                 d.get("title", ""), F_COVER_TITLE, False, C_WHITE,
                 PP_ALIGN.CENTER)
    # 分隔線
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(5.2), Inches(4.5), Inches(2.0), Inches(0.004))
    sep.fill.solid()
    sep.fill.fore_color.rgb = C_WHITE
    sep.line.fill.background()
    # 標語
    _add_textbox(slide, Inches(3.5), Inches(4.8), Inches(6.3), Inches(0.4),
                 d.get("tagline", ""), F_COVER_TAG, False, C_WHITE,
                 PP_ALIGN.CENTER)
    return slide


def slide_toc(prs, d):
    """Slide 2: 目錄"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_page_title(slide, "目錄")

    chapters_left = [
        "1. 個人介紹",
        "2. 朋友圈總覽",
        "3. 核心圈（創業夥伴與交惡對象）",
        "4. 企業圈（董事會與經營團隊）",
        "5. 協會與政商圈",
        "6. 政策與產業推手圈",
    ]
    chapters_right = [
        "7. 學術校友圈",
        "8. 思想與論壇引路人圈",
        "9. 媒體與公開平台圈",
        "10. 必須迴避對象與引薦建議",
        "11. 家族成員及企業版圖",
    ]

    fc = d.get("friend_circles", {})
    cards = fc.get("overview_cards", [])
    # 根據實際資料動態更新目錄標題
    if cards:
        for card in cards:
            cid = card.get("id", 0)
            if cid == 1 and card.get("subtitle"):
                chapters_left[2] = f"3. 核心圈（{card['subtitle']}）"
            elif cid == 2 and card.get("subtitle"):
                chapters_left[3] = f"4. 企業圈（{card['subtitle']}）"
            elif cid == 3 and card.get("subtitle"):
                chapters_left[4] = f"5. 協會與政商圈（{card['subtitle']}）"

    y = Inches(1.3)
    for ch in chapters_left:
        _add_textbox(slide, Inches(0.8), y, Inches(5.5), Inches(0.4),
                     ch, F_TOC, True, C_DEEP_BLUE)
        y += Inches(0.55)

    y = Inches(1.29)
    for ch in chapters_right:
        _add_textbox(slide, Inches(7.0), y, Inches(5.5), Inches(0.4),
                     ch, F_TOC, True, C_DEEP_BLUE)
        y += Inches(0.55)

    _add_page_num(slide, 1)
    return slide


def slide_personal(prs, d):
    """Slide 3: 個人介紹"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    p_data = d.get("personal", {})
    name = d.get("target_name", "")
    _add_page_title(slide, f"{name}簡介")

    # 左側 —— 基本資料
    left_x = Inches(0.5)
    y = Inches(1.1)
    name_en = d.get("target_name_en", "")
    full_name = f"{name} {name_en}" if name_en else name
    _add_textbox(slide, left_x, y, Inches(3.5), Inches(0.6),
                 full_name, Pt(28), True, C_DEEP_BLUE)
    y += Inches(0.5)

    # 基本資料行
    info_lines = []
    birth = p_data.get("birth_year", "")
    place = p_data.get("birth_place", "")
    age = p_data.get("age_note", "")
    if birth:
        info_lines.append(f"{birth}生（{age}）｜{place}" if place else f"{birth}生（{age}）")
    for pos in p_data.get("current_positions", []):
        info_lines.append(pos)
    sh = p_data.get("shareholding", "")
    if sh:
        info_lines.append(sh)
    interests = p_data.get("investment_interests", "")
    if interests:
        info_lines.append(f"投資興趣：{interests}")

    info_text = "\n".join(info_lines)
    _add_textbox(slide, left_x, y, Inches(4.8), Inches(1.2),
                 info_text, F_BODY, False, C_DARK_TEXT)
    y += Inches(1.3)

    # 學歷/現任/榮譽 table
    edu = p_data.get("education", [])
    positions = p_data.get("current_positions", [])
    honors = p_data.get("honors", [])

    table_data = []
    if edu:
        table_data.append(["學歷", "\n".join(e.get("school", "") + (f"（{e.get('note')}）" if e.get("note") else "") for e in edu)])
    if positions:
        table_data.append(["現任", "\n".join(positions)])
    if honors:
        table_data.append(["榮譽", "\n".join(honors)])

    if table_data:
        _make_table(slide, table_data, left_x, y, Inches(5.3),
                    [1.2, 4.5], Inches(0.55))

    # 右側 —— 職涯里程碑
    right_x = Inches(6.2)
    _add_section_title(slide, "職涯關鍵里程碑", right_x, Inches(1.1))

    milestones = p_data.get("career_milestones", [])
    my = Inches(1.5)
    for ms in milestones[:12]:  # 最多 12 個
        text = f"{ms.get('year', '')} — {ms.get('event', '')}"
        _add_textbox(slide, right_x, my, Inches(6.5), Inches(0.3),
                     text, F_MILESTONE, True, C_DEEP_BLUE)
        my += Inches(0.3)

    # 經營哲學框
    phil = p_data.get("philosophy", {})
    if phil:
        box_top = my + Inches(0.15)
        _add_rounded_rect(slide, right_x, box_top, Inches(6.5), Inches(1.3),
                          RGBColor(0xFE, 0xF3, 0xC7))
        _add_textbox(slide, Inches(6.4), box_top + Inches(0.1), Inches(6.2), Inches(0.3),
                     phil.get("title", "核心經營哲學"), F_BODY, True, C_DARK_BROWN)
        quotes_text = "\n".join(phil.get("quotes", []))
        _add_textbox(slide, Inches(6.4), box_top + Inches(0.4), Inches(6.2), Inches(0.9),
                     quotes_text, F_BODY, False, C_BROWN)

    src = p_data.get("source", "")
    _add_source_line(slide, src)
    _add_page_num(slide, 2)
    return slide


def slide_overview_dashboard(prs, d):
    """Slide 4: 朋友圈總覽 Dashboard"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fc = d.get("friend_circles", {})
    cards = fc.get("overview_cards", [])

    overview_title = "朋友圈總覽"
    # 嘗試取得核心圈 title 的副標題
    core = fc.get("core_circle", {})
    if core.get("title"):
        parts = core["title"].split("  ")
        if len(parts) > 1:
            overview_title = f"朋友圈總覽  「{parts[-1].strip('（）()')}」"

    _add_page_title(slide, overview_title)

    # 3 columns x 3 rows grid
    cols_x = [Inches(0.5), Inches(4.55), Inches(8.6)]
    rows_y = [Inches(1.2), Inches(2.7), Inches(4.2)]
    card_w = Inches(3.9)
    card_h = Inches(1.3)

    for idx, card in enumerate(cards[:9]):
        col = idx % 3
        row = idx // 3
        cx, cy = cols_x[col], rows_y[row]

        _add_rounded_rect(slide, cx, cy, card_w, card_h, C_CARD_BG)

        _add_textbox(slide, cx + Inches(0.15), cy + Inches(0.1), card_w - Inches(0.3), Inches(0.3),
                     f"{card.get('id', idx+1)}. {card.get('title', '')}", F_CARD_TITLE, True, C_GOLD)
        _add_textbox(slide, cx + Inches(0.15), cy + Inches(0.45), card_w - Inches(0.3), Inches(0.25),
                     card.get("subtitle", ""), F_CARD_SUB, False, C_LIGHT_GRAY)
        _add_textbox(slide, cx + Inches(0.15), cy + Inches(0.75), card_w - Inches(0.3), Inches(0.4),
                     card.get("people", ""), F_CARD_PEOPLE, False, C_WHITE)

    src = fc.get("core_circle", {}).get("source", "")
    if not src:
        src = "來源：今周刊、天下雜誌、工商時報、億光官網"
    _add_source_line(slide, src)
    _add_page_num(slide, 3)
    return slide


def slide_core_circle(prs, d):
    """Slide 5: 核心圈"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    cc = d.get("friend_circles", {}).get("core_circle", {})
    _add_page_title(slide, cc.get("title", "核心圈"))

    y = Inches(1.1)
    pos = cc.get("positive", [])
    neg = cc.get("negative", [])

    if pos:
        _add_textbox(slide, MARGIN_L, y, Inches(5.5), Inches(0.3),
                     "🟢 正向關係", Pt(14), True, C_GREEN)
        y += Inches(0.3)
        tdata = [["姓名", "職位", "關係", "關鍵互動脈絡"]]
        for p in pos:
            tdata.append([p.get("name",""), p.get("position",""), p.get("relation",""), p.get("detail","")])
        _make_table(slide, tdata, MARGIN_L, y, Inches(12.4), [1.2, 3, 1, 5], Inches(0.42))
        y += Inches(0.42 * len(tdata)) + Inches(0.2)

    if neg:
        _add_textbox(slide, MARGIN_L, y, Inches(5.5), Inches(0.3),
                     "🔴 負向關係（須迴避）", Pt(14), True, C_RED)
        y += Inches(0.3)
        tdata = [["姓名", "職位", "關係", "關鍵互動脈絡"]]
        for n in neg:
            tdata.append([n.get("name",""), n.get("position",""), n.get("relation",""), n.get("detail","")])
        _make_table(slide, tdata, MARGIN_L, y, Inches(12.4), [1.2, 3, 1, 5], Inches(0.42))

    _add_source_line(slide, cc.get("source", ""))
    _add_page_num(slide, 4)
    return slide


def slide_business_circle(prs, d):
    """Slide 6: 企業圈"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bc = d.get("friend_circles", {}).get("business_circle", {})
    _add_page_title(slide, bc.get("title", "企業圈"))

    # 左：董事會
    board = bc.get("board", [])
    if board:
        _add_section_title(slide, "現任董事會成員", Inches(0.5), Inches(1.1))
        tdata = [["姓名", "職位/背景", "關係"]]
        for b in board:
            tdata.append([b.get("name",""), b.get("position",""), b.get("relation","")])
        _make_table(slide, tdata, Inches(0.5), Inches(1.6), Inches(5.8), [1.5, 3.5, 1], Inches(0.38))

    # 右：經營團隊
    mgmt = bc.get("management", [])
    if mgmt:
        _add_section_title(slide, "核心經營團隊", Inches(6.6), Inches(1.1))
        tdata = [["姓名", "職位與專長"]]
        for m in mgmt:
            tdata.append([m.get("name",""), m.get("position","")])
        _make_table(slide, tdata, Inches(6.6), Inches(1.6), Inches(5.8), [1.5, 4.5], Inches(0.38))

    _add_source_line(slide, bc.get("source", ""))
    _add_page_num(slide, 5)
    return slide


def _slide_generic_table(prs, d, circle_key, page_num):
    """通用：協會/政策/學術/思想/媒體 用表格呈現的頁面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    circle = d.get("friend_circles", {}).get(circle_key, {})
    _add_page_title(slide, circle.get("title", ""))

    sections = circle.get("sections", [])
    y = Inches(1.1)

    for sec in sections:
        if y > Inches(6.0):
            break
        st = sec.get("section_title", "")
        if st:
            _add_section_title(slide, st, MARGIN_L, y)
            y += Inches(0.35)
        desc = sec.get("description", "")
        if desc:
            _add_textbox(slide, MARGIN_L, y, Inches(12), Inches(0.3),
                         desc, F_BODY, False, C_DARK_TEXT)
            y += Inches(0.35)
        people = sec.get("people", [])
        if people:
            cols_available = ["name", "position", "relation", "detail"]
            headers = []
            has_cols = set()
            for p in people:
                for k in cols_available:
                    if p.get(k):
                        has_cols.add(k)
            col_map = {"name": "姓名", "position": "職稱", "relation": "關係", "detail": "互動事證"}
            headers = [col_map[k] for k in cols_available if k in has_cols]
            tdata = [headers]
            for p in people:
                row = [p.get(k, "") for k in cols_available if k in has_cols]
                tdata.append(row)
            widths = [1.5] * len(headers)
            if len(headers) >= 4:
                widths = [1.5, 3, 1, 4]
            elif len(headers) == 3:
                widths = [1.5, 3, 1.5]
            _make_table(slide, tdata, MARGIN_L, y, Inches(12.4), widths, Inches(0.4))
            y += Inches(0.4 * len(tdata)) + Inches(0.3)

    # 如果有 people 直接在 circle 層（非 sections）
    people_direct = circle.get("people", [])
    if people_direct and not sections:
        headers_keys = []
        sample = people_direct[0] if people_direct else {}
        possible_keys = [("name","姓名"), ("position","現任職稱"), ("domain","所屬領域"),
                         ("relation","關係"), ("detail","互動事證")]
        for k, h in possible_keys:
            if sample.get(k) is not None:
                headers_keys.append((k, h))
        tdata = [[h for _, h in headers_keys]]
        for p in people_direct:
            tdata.append([p.get(k, "") for k, _ in headers_keys])
        widths = [2] * len(headers_keys)
        _make_table(slide, tdata, MARGIN_L, y, Inches(12.4), widths, Inches(0.4))
        y += Inches(0.4 * len(tdata)) + Inches(0.2)

    # early_relations 附註
    er = circle.get("early_relations", "")
    if er and y < Inches(6.5):
        _add_rounded_rect(slide, MARGIN_L, y, Inches(12.4), Inches(0.5),
                          RGBColor(0xF7, 0xFA, 0xFC))
        _add_textbox(slide, Inches(0.6), y + Inches(0.08), Inches(12), Inches(0.35),
                     er, F_BODY, True, C_DEEP_BLUE)

    _add_source_line(slide, circle.get("source", ""))
    _add_page_num(slide, page_num)
    return slide


def slide_academic(prs, d, page_num):
    """Slide 9: 學術校友圈（特殊 layout）"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    ac = d.get("friend_circles", {}).get("academic_circle", {})
    _add_page_title(slide, ac.get("title", "學術校友圈"))

    # 左：學歷背景
    _add_rounded_rect(slide, MARGIN_L, Inches(1.1), Inches(6.0), Inches(1.5),
                      RGBColor(0xF7, 0xFA, 0xFC))
    _add_textbox(slide, Inches(0.6), Inches(1.2), Inches(5.5), Inches(0.25),
                 "學歷背景", F_BODY, True, C_DEEP_BLUE)
    edu_text = "\n".join(f"• {e}" for e in ac.get("education_bg", []))
    _add_textbox(slide, Inches(0.6), Inches(1.5), Inches(5.5), Inches(1.0),
                 edu_text, F_BODY, False, C_DARK_TEXT)

    # 右：母校貢獻
    _add_rounded_rect(slide, Inches(6.7), Inches(1.1), Inches(6.0), Inches(1.5),
                      RGBColor(0xFE, 0xF3, 0xC7))
    _add_textbox(slide, Inches(6.9), Inches(1.2), Inches(5.5), Inches(0.25),
                 "母校貢獻與榮譽", F_BODY, True, C_DARK_BROWN)
    contrib_text = "\n".join(f"• {c}" for c in ac.get("school_contribution", []))
    _add_textbox(slide, Inches(6.9), Inches(1.5), Inches(5.5), Inches(1.0),
                 contrib_text, F_BODY, False, C_BROWN)

    # 校友圈表格
    alumni = ac.get("alumni", [])
    if alumni:
        _add_section_title(slide, "校友圈人物", MARGIN_L, Inches(2.8))
        tdata = [["校友名稱", "現任職稱", "關係", "互動證據"]]
        for a in alumni:
            tdata.append([a.get("name",""), a.get("position",""), a.get("relation",""), a.get("detail","")])
        _make_table(slide, tdata, MARGIN_L, Inches(3.2), Inches(12.4), [1.5, 3, 1, 4.5], Inches(0.4))

    # 底部引言
    note = ac.get("extra_note", "")
    if note:
        _add_rounded_rect(slide, MARGIN_L, Inches(5.6), Inches(12.4), Inches(0.5),
                          RGBColor(0xF7, 0xFA, 0xFC))
        _add_textbox(slide, Inches(0.6), Inches(5.65), Inches(12), Inches(0.4),
                     note, F_BODY, False, C_DARK_TEXT)

    _add_source_line(slide, ac.get("source", ""))
    _add_page_num(slide, page_num)
    return slide


def slide_thought_leaders(prs, d, page_num):
    """Slide 10: 思想與論壇引路人圈"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tl = d.get("friend_circles", {}).get("thought_leaders", {})
    _add_page_title(slide, tl.get("title", "思想與論壇引路人圈"))

    # 思想影響者表格
    _add_section_title(slide, "思想影響者", MARGIN_L, Inches(1.1))
    people = tl.get("people", [])
    if people:
        tdata = [["姓名", "身份/職稱", "關係", "影響說明"]]
        for p in people:
            tdata.append([p.get("name",""), p.get("position",""), p.get("relation",""), p.get("detail","")])
        _make_table(slide, tdata, MARGIN_L, Inches(1.5), Inches(12.4), [1.5, 3, 1, 5], Inches(0.42))

    # 經典語錄
    quotes = tl.get("quotes", [])
    if quotes:
        name = d.get("target_name", "")
        _add_section_title(slide, f"{name}經典語錄（反映經營哲學）", MARGIN_L, Inches(3.7),
                           Inches(8))
        # 修改顏色為金色
        qy = Inches(4.1)
        for q in quotes[:4]:
            _add_rounded_rect(slide, MARGIN_L, qy, Inches(12.4), Inches(0.45),
                              RGBColor(0xFE, 0xF3, 0xC7))
            _add_textbox(slide, Inches(0.6), qy + Inches(0.08), Inches(12), Inches(0.35),
                         q, F_QUOTE, True, C_BROWN)
            qy += Inches(0.55)

    _add_source_line(slide, tl.get("source", ""))
    _add_page_num(slide, page_num)
    return slide


def slide_media(prs, d, page_num):
    """Slide 11: 媒體與公開平台圈"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    mc = d.get("friend_circles", {}).get("media_circle", {})
    _add_page_title(slide, mc.get("title", "媒體與公開平台圈"))

    y = Inches(1.1)
    pos_events = mc.get("positive_events", [])
    neg_events = mc.get("negative_events", [])

    if pos_events:
        _add_textbox(slide, MARGIN_L, y, Inches(5.5), Inches(0.3),
                     "🟢 正面事件", Pt(14), True, C_GREEN)
        y += Inches(0.3)
        tdata = [["時間", "媒體/平台", "性質", "內容摘要"]]
        for e in pos_events:
            tdata.append([e.get("time",""), e.get("media",""), e.get("nature",""), e.get("summary","")])
        _make_table(slide, tdata, MARGIN_L, y, Inches(12.4), [1.2, 2, 1, 5.5], Inches(0.38))
        y += Inches(0.38 * len(tdata)) + Inches(0.2)

    if neg_events:
        _add_textbox(slide, MARGIN_L, y, Inches(5.5), Inches(0.3),
                     "🔴 負面事件", Pt(14), True, C_RED)
        y += Inches(0.3)
        tdata = [["時間", "媒體/平台", "性質", "內容摘要"]]
        for e in neg_events:
            tdata.append([e.get("time",""), e.get("media",""), e.get("nature",""), e.get("summary","")])
        _make_table(slide, tdata, MARGIN_L, y, Inches(12.4), [1.2, 2, 1, 5.5], Inches(0.38))

    _add_source_line(slide, mc.get("source", ""))
    _add_page_num(slide, page_num)
    return slide


def slide_avoidance(prs, d, page_num):
    """Slide 12: 迴避對象與引薦建議"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    av = d.get("friend_circles", {}).get("avoidance", {})
    _add_page_title(slide, av.get("title", "必須迴避對象與引薦建議"))

    # 左：迴避
    _add_textbox(slide, MARGIN_L, Inches(1.1), Inches(5.5), Inches(0.3),
                 "❌ 必須迴避對象", Pt(14), True, C_RED)
    avoid = av.get("avoid_list", [])
    if avoid:
        tdata = [["姓名", "現職", "迴避原因"]]
        for a in avoid:
            tdata.append([a.get("name",""), a.get("position",""), a.get("reason","")])
        _make_table(slide, tdata, MARGIN_L, Inches(1.4), Inches(6.0), [2, 2, 3], Inches(0.38))

    # 右：引薦
    _add_textbox(slide, Inches(6.7), Inches(1.1), Inches(5.5), Inches(0.3),
                 "✅ 優先引薦管道", Pt(14), True, C_GREEN)
    rec = av.get("recommend_list", [])
    if rec:
        tdata = [["姓名", "現職", "引薦方式"]]
        for r in rec:
            tdata.append([r.get("name",""), r.get("position",""), r.get("method","")])
        _make_table(slide, tdata, Inches(6.7), Inches(1.4), Inches(5.9), [2, 2, 3], Inches(0.38))

    _add_source_line(slide, av.get("source", ""))
    _add_page_num(slide, page_num)
    return slide


def slide_family_tree(prs, d, page_num):
    """
    Slide 13: 家族樹 — 自動生成
    從 family.tree_data 讀取世代資料，用 shapes + connector lines 畫家族樹。
    tree_data 格式：
    [
      {"generation": "G0", "members": [
        {"name":"祖父","relation":"祖父","title":"...","gender":"M"},
        {"name":"祖母","relation":"祖母","title":"...","gender":"F"}
      ]},
      {"generation": "G1", "members": [...]},
      ...
    ]
    如果 family_tree_image_path 存在且檔案存在，優先用圖片。
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_page_title(slide, "家族樹")

    # 優先用圖片
    img_path = d.get("family", {}).get("family_tree_image_path")
    if img_path and os.path.exists(str(img_path)):
        slide.shapes.add_picture(str(img_path), Inches(1.5), Inches(1.0),
                                 Inches(10.2), Inches(5.8))
        _add_page_num(slide, page_num)
        return slide

    # 自動生成家族樹
    tree_data = d.get("family", {}).get("tree_data", [])
    if not tree_data:
        # 如果沒有 tree_data，嘗試從 members_detail 推斷
        _add_textbox(slide, Inches(2), Inches(3.5), Inches(9), Inches(0.5),
                     "（tree_data 為空，無法自動生成家族樹。請在 JSON 中填入 family.tree_data）",
                     F_BODY, False, C_GRAY, PP_ALIGN.CENTER)
        _add_page_num(slide, page_num)
        return slide

    # Layout 計算
    BOX_W = Inches(1.55)
    BOX_H = Inches(0.85)
    H_GAP = Inches(0.3)
    V_GAP = Inches(0.45)
    start_y = Inches(1.15)
    page_center = Inches(6.666)  # 13.333 / 2

    # 顏色定義
    gen_colors = {
        "G0": RGBColor(0xBE, 0xE3, 0xF8),  # 淺藍
        "G1": RGBColor(0xC6, 0xF6, 0xD5),  # 淺綠
        "G2": RGBColor(0xFE, 0xF3, 0xC7),  # 淺黃（主角世代）
        "G3": RGBColor(0xFE, 0xD7, 0xD7),  # 淺紅
    }
    target_name = d.get("target_name", "")

    all_boxes = []  # [(gen_idx, member_idx, center_x, center_y, name)]

    for gen_idx, gen in enumerate(tree_data[:4]):  # 最多4代
        gen_label = gen.get("generation", f"G{gen_idx}")
        members = gen.get("members", [])
        if not members:
            continue

        y = start_y + gen_idx * (BOX_H + V_GAP)
        n = len(members)
        total_w = n * BOX_W + (n - 1) * H_GAP
        start_x = page_center - total_w / 2

        # 世代標籤
        label_tb = slide.shapes.add_textbox(
            Inches(0.2), y + Inches(0.15), Inches(0.7), Inches(0.4)
        )
        lp = label_tb.text_frame.paragraphs[0]
        lp.text = gen_label
        lp.font.size = Pt(10)
        lp.font.bold = True
        lp.font.color.rgb = C_GRAY

        fill_color = gen_colors.get(gen_label, C_HEADER_BG)

        for m_idx, member in enumerate(members):
            x = start_x + m_idx * (BOX_W + H_GAP)
            name = member.get("name", "N/A")
            relation = member.get("relation", "")
            title = member.get("title", "")

            # 主角用深藍邊框
            is_target = (name == target_name)
            border_color = C_DEEP_BLUE if is_target else C_GRAY

            # Box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, int(x), int(y),
                int(BOX_W), int(BOX_H)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = fill_color
            box.line.color.rgb = border_color
            box.line.width = Pt(2) if is_target else Pt(1)

            # 文字
            tf = box.text_frame
            tf.word_wrap = True
            tf.margin_left = Inches(0.05)
            tf.margin_right = Inches(0.05)
            tf.margin_top = Inches(0.03)
            tf.margin_bottom = Inches(0.03)

            p1 = tf.paragraphs[0]
            p1.text = name
            p1.font.size = Pt(10)
            p1.font.bold = True
            p1.font.color.rgb = C_DEEP_BLUE if is_target else C_DARK_TEXT
            p1.alignment = PP_ALIGN.CENTER

            p2 = tf.add_paragraph()
            p2.text = relation
            p2.font.size = Pt(8)
            p2.font.color.rgb = C_GRAY
            p2.alignment = PP_ALIGN.CENTER

            if title:
                p3 = tf.add_paragraph()
                display_title = title[:18] + "..." if len(title) > 18 else title
                p3.text = display_title
                p3.font.size = Pt(7)
                p3.font.color.rgb = C_GRAY
                p3.alignment = PP_ALIGN.CENTER

            cx = x + BOX_W / 2
            cy_top = y
            cy_bot = y + BOX_H
            all_boxes.append((gen_idx, m_idx, int(cx), int(cy_top), int(cy_bot), name))

    # 畫世代間的連接線（簡化版：每代中心往下一代中心畫線）
    from pptx.oxml.ns import qn
    for gen_idx in range(len(tree_data) - 1):
        parents = [b for b in all_boxes if b[0] == gen_idx]
        children = [b for b in all_boxes if b[0] == gen_idx + 1]
        if not parents or not children:
            continue

        # 計算 parent 群的中心 x
        p_cx = sum(b[2] for b in parents) // len(parents)
        p_bot = parents[0][4]  # bottom y of parent row

        # 畫一條從 parent 底部往下的垂直線
        mid_y = p_bot + int(V_GAP * 0.45)
        connector_v = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            p_cx - Inches(0.005), p_bot,
            Inches(0.01), mid_y - p_bot
        )
        connector_v.fill.solid()
        connector_v.fill.fore_color.rgb = C_GRAY
        connector_v.line.fill.background()

        # 水平分支線
        if len(children) > 1:
            c_left = min(b[2] for b in children)
            c_right = max(b[2] for b in children)
            connector_h = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                c_left, mid_y - Inches(0.005),
                c_right - c_left, Inches(0.01)
            )
            connector_h.fill.solid()
            connector_h.fill.fore_color.rgb = C_GRAY
            connector_h.line.fill.background()

        # 每個 child 的垂直線
        for child in children:
            c_cx = child[2]
            c_top = child[3]
            cv = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                c_cx - Inches(0.005), mid_y,
                Inches(0.01), c_top - mid_y
            )
            cv.fill.solid()
            cv.fill.fore_color.rgb = C_GRAY
            cv.line.fill.background()

    _add_page_num(slide, page_num)
    return slide


def slides_member_detail(prs, d, start_page):
    """Slide 14-16: 核心成員介紹"""
    slides = []
    members = d.get("family", {}).get("members_detail", [])
    pn = start_page

    for member_slide in members:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slides.append(slide)
        _add_page_title(slide, member_slide.get("slide_title", "核心成員介紹"))

        y = Inches(1.1)
        profiles = member_slide.get("profiles", [])
        for profile in profiles:
            if y > Inches(6.0):
                break
            sub = profile.get("subtitle", "")
            if sub:
                _add_textbox(slide, MARGIN_L, y, Inches(10), Inches(0.35),
                             sub, F_SECTION, True, C_MID_BLUE)
                y += Inches(0.4)

            tbl = profile.get("table_data", [])
            if tbl:
                _make_table(slide, tbl, Inches(0.5), y, Inches(5.8), [1.5, 4.5], Inches(0.35))
                y += Inches(0.35 * len(tbl)) + Inches(0.15)

            achievements = profile.get("achievements", [])
            if achievements:
                ach_title = profile.get("achievements_title", "重要成就")
                _add_textbox(slide, MARGIN_L, y, Inches(10), Inches(0.3),
                             ach_title, F_BODY, True, C_DEEP_BLUE)
                y += Inches(0.3)
                ach_text = "\n".join(f"• {a}" for a in achievements)
                _add_textbox(slide, Inches(0.5), y, Inches(11.5), Inches(1.2),
                             ach_text, F_BODY, False, C_DARK_TEXT)
                y += Inches(0.3 * len(achievements))

        src = member_slide.get("source", "")
        _add_source_line(slide, src)
        _add_page_num(slide, pn)
        pn += 1

    return slides, pn


def slide_business_map(prs, d, page_num):
    """Slide 17: 家族企業版圖"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bm = d.get("family", {}).get("business_map", {})
    _add_page_title(slide, bm.get("title", "家族企業版圖"))

    y = Inches(1.1)
    # 家族成員任職總覽
    roles = bm.get("family_roles", [])
    if roles:
        _add_section_title(slide, "家族關聯企業與任職總覽", MARGIN_L, y)
        y += Inches(0.4)
        tdata = [["成員", "與主角關係", "職務", "參與程度"]]
        for r in roles:
            tdata.append([r.get("member",""), r.get("relation",""), r.get("position",""), r.get("involvement","")])
        _make_table(slide, tdata, MARGIN_L, y, Inches(12.4), [1.5, 1.5, 3, 4], Inches(0.38))
        y += Inches(0.38 * len(tdata)) + Inches(0.3)

    # 投資控股公司
    inv = bm.get("investment_companies", [])
    if inv:
        _add_section_title(slide, "家族關聯投資控股公司", MARGIN_L, y)
        y += Inches(0.4)
        tdata = [["投資公司名稱", "代表人/關聯人", "說明"]]
        for i in inv:
            tdata.append([i.get("company",""), i.get("representative",""), i.get("note","")])
        _make_table(slide, tdata, MARGIN_L, y, Inches(12.4), [3, 2, 5], Inches(0.38))

    _add_source_line(slide, bm.get("source", ""))
    _add_page_num(slide, page_num)
    return slide


def slide_summary(prs, d, page_num):
    """Slide 18: 總結"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    sm = d.get("summary", {})
    _add_page_title(slide, sm.get("title", "總結"))

    # 三欄
    cols = [
        ("✓ 正面觀察", sm.get("positive", []), C_GREEN),
        ("⚠ 風險關注", sm.get("risks", []), RGBColor(0xD6, 0x9E, 0x2E)),
        ("➤ 建議追蹤", sm.get("tracking", []), C_MID_BLUE),
    ]
    cx_list = [Inches(0.5), Inches(4.6), Inches(8.7)]
    col_w = Inches(3.8)

    for idx, (header, items, color) in enumerate(cols):
        cx = cx_list[idx]
        _add_textbox(slide, cx, Inches(1.1), col_w, Inches(0.35),
                     header, F_SECTION, True, color)
        items_text = "\n".join(f"• {it}" for it in items)
        _add_textbox(slide, cx, Inches(1.5), col_w, Inches(2.5),
                     items_text, F_BODY, False, C_DARK_TEXT)

    # 關鍵結論
    conclusion = sm.get("conclusion", "")
    if conclusion:
        _add_textbox(slide, MARGIN_L, Inches(4.3), Inches(4), Inches(0.35),
                     "關鍵結論", F_SECTION, True, C_DEEP_BLUE)
        _add_textbox(slide, MARGIN_L, Inches(4.7), Inches(12.4), Inches(1.5),
                     conclusion, F_BODY, False, C_DARK_TEXT)

    _add_source_line(slide, sm.get("source", ""))
    _add_page_num(slide, page_num)
    return slide


# ==============================================================================
# Main
# ==============================================================================

def generate_report(data, output_path):
    """生成完整的 18 頁人脈網絡報告"""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Slide 1: 封面
    slide_cover(prs, data)

    # Slide 2: 目錄
    slide_toc(prs, data)

    # Slide 3: 個人介紹
    slide_personal(prs, data)

    # Slide 4: 朋友圈總覽
    slide_overview_dashboard(prs, data)

    # Slide 5: 核心圈
    slide_core_circle(prs, data)

    # Slide 6: 企業圈
    slide_business_circle(prs, data)

    # Slide 7: 協會與政商圈
    _slide_generic_table(prs, data, "association_circle", 6)

    # Slide 8: 政策與產業推手圈
    _slide_generic_table(prs, data, "policy_circle", 7)

    # Slide 9: 學術校友圈
    slide_academic(prs, data, 8)

    # Slide 10: 思想與論壇引路人圈
    slide_thought_leaders(prs, data, 9)

    # Slide 11: 媒體與公開平台圈
    slide_media(prs, data, 10)

    # Slide 12: 迴避對象+引薦
    slide_avoidance(prs, data, 11)

    # Slide 13: 家族樹
    slide_family_tree(prs, data, 12)

    # Slide 14-16: 核心成員
    member_slides, next_pn = slides_member_detail(prs, data, 13)

    # Slide 17: 家族企業版圖
    slide_business_map(prs, data, next_pn)

    # Slide 18: 總結
    slide_summary(prs, data, next_pn + 1)

    prs.save(output_path)
    print(f"✅ 報告已生成：{output_path}")
    print(f"   共 {len(prs.slides)} 頁")
    return output_path


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="人脈網絡報告 PPT Generator")
    parser.add_argument("--input", "-i", required=True, help="Input JSON file")
    parser.add_argument("--output", "-o", default="/home/claude/output.pptx", help="Output PPTX path")
    args = parser.parse_args()

    with open(args.input, "r", encoding="utf-8") as f:
        data = json.load(f)
    generate_report(data, args.output)
